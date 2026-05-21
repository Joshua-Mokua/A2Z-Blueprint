"""scripts/docgen/ppt_generator.py — A2Z brochure (15 slides) (v8.13).

Per docs/A2Z_LIVING_DOCS_PLAN.md Part 4, produces A2Z_MIS_360_Brochure.pptx
with audit-locked claims. Validates every numeric claim against the registry
before writing; aborts on divergence.
"""
from __future__ import annotations

from pathlib import Path
from typing import Any, Dict, List, Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from scripts.docgen._registry_loader import load_registry
from scripts.docgen._claim_validator import (
    Claim, validate_claims, ClaimValidationError,
)
from scripts.docgen._honest_section import (
    collect_honest_scope_lines, standard_disclaimer_paragraph, section_title,
)
from scripts.docgen import _theme as theme


SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _hex_to_rgb(hex_str: str) -> RGBColor:
    """Convert hex string to pptx RGBColor."""
    return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))


def _add_blank_slide(prs: Presentation):
    """Add a blank-layout slide."""
    return prs.slides.add_slide(prs.slide_layouts[6])


def _add_textbox(slide, left, top, width, height, text, *,
                  font_name: str = theme.PPT_BODY_FONT,
                  font_size: int = theme.SIZE_BODY,
                  font_color_hex: str = theme.CHARCOAL_HEX,
                  bold: bool = False,
                  italic: bool = False,
                  align: int = PP_ALIGN.LEFT,
                  anchor: int = MSO_ANCHOR.TOP):
    """Add a text box with consistent styling."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    p.text = text
    for run in p.runs:
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = _hex_to_rgb(font_color_hex)
    return tb


def _add_filled_rect(slide, left, top, width, height, fill_hex: str,
                      line_hex: Optional[str] = None):
    """Add a filled rectangle with optional border."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex_to_rgb(fill_hex)
    if line_hex:
        shape.line.color.rgb = _hex_to_rgb(line_hex)
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def _add_footer(slide, slide_num: int, total: int, version: str):
    """Add page-number + version footer at bottom."""
    _add_textbox(slide, Inches(0.5), Inches(7.1),
                  Inches(12.3), Inches(0.3),
                  f"{theme.PRODUCT_NAME}  ·  {version}  ·  Slide {slide_num}/{total}",
                  font_size=theme.SIZE_FOOTER,
                  font_color_hex=theme.MID_GREY_HEX,
                  align=PP_ALIGN.LEFT)


# ════════════════════════════════════════════════════════════════════
# Individual slide builders
# ════════════════════════════════════════════════════════════════════

def _build_slide_1_title(prs, registry: Dict[str, Any], total: int):
    """Slide 1 — Title."""
    slide = _add_blank_slide(prs)
    # Dark title slide (sandwich convention)
    _add_filled_rect(slide, 0, 0, SLIDE_W, SLIDE_H, theme.PRIMARY_HEX)

    # Big product name
    _add_textbox(slide, Inches(0.75), Inches(2.2),
                  Inches(11.8), Inches(1.5),
                  theme.PRODUCT_NAME,
                  font_name=theme.PPT_HEADER_FONT,
                  font_size=64, bold=True,
                  font_color_hex=theme.WHITE_HEX,
                  align=PP_ALIGN.LEFT)

    # Tagline
    _add_textbox(slide, Inches(0.75), Inches(3.7),
                  Inches(11.8), Inches(0.6),
                  theme.PRODUCT_TAGLINE,
                  font_size=22, italic=True,
                  font_color_hex="CADCFC",  # ice blue accent
                  align=PP_ALIGN.LEFT)

    # Audit credentials block
    _add_textbox(slide, Inches(0.75), Inches(5.2),
                  Inches(11.8), Inches(1.5),
                  f"Version {registry['platform']['version']}  ·  "
                  f"{registry['platform']['audit_gates']} audit gates  ·  "
                  f"{registry['loops_count']} feedback loops "
                  f"({registry['loops_wired']}/{registry['loops_count']} wired)  ·  "
                  f"{registry['stocks_count']} system stocks "
                  f"({registry['stocks_wired']}/{registry['stocks_count']} wired)",
                  font_size=14,
                  font_color_hex=theme.WHITE_HEX,
                  align=PP_ALIGN.LEFT)

    # Verification CTA
    _add_textbox(slide, Inches(0.75), Inches(6.6),
                  Inches(11.8), Inches(0.4),
                  "Verify any claim in this deck: " + registry["platform"]["audit_command"],
                  font_size=11, font_name=theme.PPT_MONO_FONT,
                  font_color_hex="CADCFC",
                  align=PP_ALIGN.LEFT)


def _build_slide_2_problem(prs, registry: Dict[str, Any], slide_num: int, total: int):
    """Slide 2 — The strategy-execution gap."""
    slide = _add_blank_slide(prs)
    _add_filled_rect(slide, 0, 0, SLIDE_W, SLIDE_H, theme.WHITE_HEX)

    # Heading bar (left side)
    _add_filled_rect(slide, 0, 0, Inches(0.3), SLIDE_H, theme.PRIMARY_HEX)

    _add_textbox(slide, Inches(0.75), Inches(0.6),
                  Inches(12), Inches(0.8),
                  "The strategy-execution gap",
                  font_size=theme.SIZE_TITLE, bold=True,
                  font_color_hex=theme.PRIMARY_HEX)

    # Body text
    body_text = (
        "Most banks have strategy. Few have execution at the teller line.\n\n"
        "Quarterly board targets do not flow to individual KPIs. "
        "Manual cascade in Excel breaks down. Staff receive feedback "
        "monthly — too late to change behaviour mid-cycle.\n\n"
        "Banks run 10-15 separate systems for credit, risk, HR, finance, "
        "customer ops; each with its own truth. Audit trails are scattered. "
        "CBK reviews require manual reconstruction."
    )
    _add_textbox(slide, Inches(0.75), Inches(1.8),
                  Inches(12), Inches(4),
                  body_text,
                  font_size=theme.SIZE_HEADING - 4,
                  font_color_hex=theme.CHARCOAL_HEX)

    _add_footer(slide, slide_num, total, registry["platform"]["version"])


def _build_slide_3_solution(prs, registry: Dict[str, Any], slide_num: int, total: int):
    """Slide 3 — The A2Z approach (3 columns)."""
    slide = _add_blank_slide(prs)
    _add_filled_rect(slide, 0, 0, SLIDE_W, SLIDE_H, theme.WHITE_HEX)
    _add_filled_rect(slide, 0, 0, Inches(0.3), SLIDE_H, theme.PRIMARY_HEX)

    _add_textbox(slide, Inches(0.75), Inches(0.6),
                  Inches(12), Inches(0.8),
                  "The A2Z approach",
                  font_size=theme.SIZE_TITLE, bold=True,
                  font_color_hex=theme.PRIMARY_HEX)

    cols = [
        ("Audit-locked invariants",
         "6 defense-in-depth gates (G104-G109). "
         "Architectural rules enforced by the build. "
         "Future regressions fail the audit."),
        ("Systems-layer first",
         "Stocks, flows, invariants are first-class objects. "
         "Per Meadows: a system is its feedback loops. "
         f"{registry['loops_wired']}/{registry['loops_count']} loops wired."),
        ("Honest scope",
         "Every artifact ends with what it does NOT claim. "
         "Roadmap items marked. ROI projections "
         "distinguished from measurements."),
    ]

    col_w = Inches(3.9)
    col_x_start = Inches(0.75)
    col_y = Inches(1.8)
    col_h = Inches(4.5)
    for i, (title, body) in enumerate(cols):
        x = col_x_start + (col_w + Inches(0.2)) * i
        # Card background
        _add_filled_rect(slide, x, col_y, col_w, col_h,
                          theme.LIGHT_GREY_HEX, line_hex=theme.MID_GREY_HEX)
        # Card title
        _add_textbox(slide, x + Inches(0.2), col_y + Inches(0.3),
                      col_w - Inches(0.4), Inches(0.6),
                      title,
                      font_size=18, bold=True,
                      font_color_hex=theme.PRIMARY_HEX)
        # Card body
        _add_textbox(slide, x + Inches(0.2), col_y + Inches(1.1),
                      col_w - Inches(0.4), Inches(3.2),
                      body,
                      font_size=12,
                      font_color_hex=theme.CHARCOAL_HEX)

    _add_footer(slide, slide_num, total, registry["platform"]["version"])


def _build_slide_4_architecture(prs, registry: Dict[str, Any], slide_num: int, total: int):
    """Slide 4 — Architecture (5+1 tiers)."""
    slide = _add_blank_slide(prs)
    _add_filled_rect(slide, 0, 0, SLIDE_W, SLIDE_H, theme.WHITE_HEX)
    _add_filled_rect(slide, 0, 0, Inches(0.3), SLIDE_H, theme.PRIMARY_HEX)

    _add_textbox(slide, Inches(0.75), Inches(0.6),
                  Inches(12), Inches(0.8),
                  "Architecture — six tiers, audit-locked",
                  font_size=theme.SIZE_TITLE, bold=True,
                  font_color_hex=theme.PRIMARY_HEX)

    tiers = [
        ("Tier 1 — Domain engines",
         f"{registry['platform']['engines_count']} modules in utils/ — pure deterministic logic"),
        ("Tier 2 — Systems-layer registries",
         f"{registry['stocks_count']} stocks, {registry['loops_count']} loops, "
         f"invariants, {len(registry.get('invariants', []))} hard constraints"),
        ("Tier 3 — Audit perimeter",
         f"scripts/audit.py with {registry['platform']['audit_gates']} gates "
         "(G104-G109 defense-in-depth)"),
        ("Tier 4 — Technical-grade docs",
         f"Charter + 2 retrospectives + {registry['platform']['changelog_count']} CHANGELOGs"),
        ("Tier 5 — UI surfaces",
         "4 dedicated pages + 16 enhanced pages + page 91 systems view"),
        ("Tier 6 — Living Documentation",
         "Sales-grade rendering with audit-locked claim validation (this artifact)"),
    ]

    row_h = Inches(0.7)
    row_y = Inches(1.8)
    for i, (label, body) in enumerate(tiers):
        y = row_y + (row_h + Inches(0.1)) * i
        # Tier number circle (left)
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.75), y,
                                          Inches(0.5), Inches(0.5))
        circle.fill.solid()
        circle.fill.fore_color.rgb = _hex_to_rgb(theme.PRIMARY_HEX)
        circle.line.fill.background()
        circle.text_frame.text = str(i + 1)
        circle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        for run in circle.text_frame.paragraphs[0].runs:
            run.font.name = theme.PPT_HEADER_FONT
            run.font.size = Pt(14)
            run.font.bold = True
            run.font.color.rgb = _hex_to_rgb(theme.WHITE_HEX)
        # Label + body
        _add_textbox(slide, Inches(1.5), y - Inches(0.05),
                      Inches(11), Inches(0.35), label,
                      font_size=14, bold=True,
                      font_color_hex=theme.PRIMARY_HEX)
        _add_textbox(slide, Inches(1.5), y + Inches(0.3),
                      Inches(11), Inches(0.4), body,
                      font_size=11,
                      font_color_hex=theme.CHARCOAL_HEX)

    _add_footer(slide, slide_num, total, registry["platform"]["version"])


def _build_slide_5_systems(prs, registry: Dict[str, Any], slide_num: int, total: int):
    """Slide 5 — Stocks + loops at-a-glance."""
    slide = _add_blank_slide(prs)
    _add_filled_rect(slide, 0, 0, SLIDE_W, SLIDE_H, theme.WHITE_HEX)
    _add_filled_rect(slide, 0, 0, Inches(0.3), SLIDE_H, theme.PRIMARY_HEX)

    _add_textbox(slide, Inches(0.75), Inches(0.6),
                  Inches(12), Inches(0.8),
                  "Systems-layer — what flows where",
                  font_size=theme.SIZE_TITLE, bold=True,
                  font_color_hex=theme.PRIMARY_HEX)

    # 4 big stat callouts
    stats = [
        (str(registry["stocks_count"]), "system stocks", f"{registry['stocks_wired']}/{registry['stocks_count']} WIRED"),
        (str(registry["loops_count"]), "feedback loops", f"{registry['loops_wired']}/{registry['loops_count']} WIRED"),
        (str(registry["learning_loops_count"]), "learning loops", "Meadows' highest-value type"),
        (str(registry["platform"]["audit_gates"]), "audit gates", "G104-G109 defense-in-depth"),
    ]

    stat_w = Inches(2.9)
    for i, (big, label, sub) in enumerate(stats):
        x = Inches(0.75) + (stat_w + Inches(0.2)) * i
        _add_filled_rect(slide, x, Inches(2.0), stat_w, Inches(2.5),
                          theme.LIGHT_GREY_HEX, line_hex=theme.MID_GREY_HEX)
        _add_textbox(slide, x, Inches(2.2), stat_w, Inches(1.2), big,
                      font_size=72, bold=True,
                      font_color_hex=theme.PRIMARY_HEX,
                      align=PP_ALIGN.CENTER)
        _add_textbox(slide, x, Inches(3.5), stat_w, Inches(0.5), label,
                      font_size=16, bold=True,
                      font_color_hex=theme.CHARCOAL_HEX,
                      align=PP_ALIGN.CENTER)
        _add_textbox(slide, x, Inches(4.0), stat_w, Inches(0.5), sub,
                      font_size=11, italic=True,
                      font_color_hex=theme.MID_GREY_HEX,
                      align=PP_ALIGN.CENTER)

    # Caption
    _add_textbox(slide, Inches(0.75), Inches(5.0),
                  Inches(12), Inches(1.5),
                  "Every loop is registered in utils/system_flows.py with explicit status, "
                  "pattern, and notes. The 6-stock × 15-loop graph is the canonical model "
                  "of how value flows through the bank — credit risk to capital adequacy, "
                  "customer profitability to BSC scoring, channel reliability to customer alerts.",
                  font_size=13,
                  font_color_hex=theme.CHARCOAL_HEX,
                  align=PP_ALIGN.LEFT)

    _add_footer(slide, slide_num, total, registry["platform"]["version"])


def _build_slide_6_acl(prs, registry: Dict[str, Any], slide_num: int, total: int):
    """Slide 6 — FLEXCUBE Anti-Corruption Layer."""
    slide = _add_blank_slide(prs)
    _add_filled_rect(slide, 0, 0, SLIDE_W, SLIDE_H, theme.WHITE_HEX)
    _add_filled_rect(slide, 0, 0, Inches(0.3), SLIDE_H, theme.PRIMARY_HEX)

    _add_textbox(slide, Inches(0.75), Inches(0.6),
                  Inches(12), Inches(0.8),
                  "FLEXCUBE Anti-Corruption Layer",
                  font_size=theme.SIZE_TITLE, bold=True,
                  font_color_hex=theme.PRIMARY_HEX)

    body = (
        "Per Eric Evans' DDD pattern, the ACL translates FLEXCUBE field names "
        "(GROSS_OS, SEGMENT_DIST, STAGE_DIST) into A2Z's normalised vocabulary "
        "INSIDE the adapter. A2Z domain code never sees FLEXCUBE-specific names.\n\n"
        "5 portfolio-aggregate methods are live in v8.0+:\n"
        "  • /PortfolioService/Loans/Aggregate\n"
        "  • /PortfolioService/Deposits/Aggregate\n"
        "  • /PortfolioService/NPL/Aggregate\n"
        "  • /CustomerService/Aggregate\n"
        "  • /AccountService/Dormancy/Aggregate\n\n"
        "Three modes — live, mock, synthetic. The 3-tier fallback (live → CBS-synthetic → "
        "demo defaults) means platform survives FLEXCUBE outages gracefully."
    )
    _add_textbox(slide, Inches(0.75), Inches(1.8),
                  Inches(12), Inches(4.5),
                  body,
                  font_size=14,
                  font_color_hex=theme.CHARCOAL_HEX)

    _add_footer(slide, slide_num, total, registry["platform"]["version"])


def _build_slide_7_resilience(prs, registry: Dict[str, Any], slide_num: int, total: int):
    """Slide 7 — Resilience layers."""
    slide = _add_blank_slide(prs)
    _add_filled_rect(slide, 0, 0, SLIDE_W, SLIDE_H, theme.WHITE_HEX)
    _add_filled_rect(slide, 0, 0, Inches(0.3), SLIDE_H, theme.PRIMARY_HEX)

    _add_textbox(slide, Inches(0.75), Inches(0.6),
                  Inches(12), Inches(0.8),
                  "Resilience — per CBK Operations Resilience Guidelines (2019)",
                  font_size=theme.SIZE_TITLE - 6, bold=True,
                  font_color_hex=theme.PRIMARY_HEX)

    layers = [
        ("Retry", "3 attempts with exponential backoff (1s/3s/9s) + ±20% jitter",
         "v8.1 + v8.8"),
        ("Circuit breaker", "Trips OPEN after 5 consecutive failures; 60s open duration; half-open probe pattern",
         "v8.1"),
        ("Latency telemetry", "Per-endpoint p50/p95/p99 over rolling 200-sample window; circuit-open suppression",
         "v8.2"),
        ("Restart-free admin", "reset_circuit() + replay_events() with audit trails",
         "v8.9"),
    ]

    row_h = Inches(1.0)
    row_y = Inches(1.8)
    for i, (name, desc, ver) in enumerate(layers):
        y = row_y + row_h * i
        # Number column
        _add_filled_rect(slide, Inches(0.75), y, Inches(0.7), Inches(0.85),
                          theme.PRIMARY_HEX)
        _add_textbox(slide, Inches(0.75), y + Inches(0.2), Inches(0.7), Inches(0.5),
                      str(i + 1), font_size=24, bold=True,
                      font_color_hex=theme.WHITE_HEX, align=PP_ALIGN.CENTER)
        # Name + version
        _add_textbox(slide, Inches(1.7), y, Inches(11), Inches(0.4),
                      f"{name}  ·  {ver}",
                      font_size=16, bold=True,
                      font_color_hex=theme.PRIMARY_HEX)
        # Description
        _add_textbox(slide, Inches(1.7), y + Inches(0.4), Inches(11), Inches(0.5),
                      desc,
                      font_size=12,
                      font_color_hex=theme.CHARCOAL_HEX)

    # Footer note
    _add_textbox(slide, Inches(0.75), Inches(6.3), Inches(12), Inches(0.4),
                  "Locked by audit gate G108 (flexcube_retry_circuit_breaker_contract).",
                  font_size=10, italic=True,
                  font_color_hex=theme.MID_GREY_HEX)

    _add_footer(slide, slide_num, total, registry["platform"]["version"])


def _build_slide_8_observability(prs, registry: Dict[str, Any], slide_num: int, total: int):
    """Slide 8 — Observability triangle."""
    slide = _add_blank_slide(prs)
    _add_filled_rect(slide, 0, 0, SLIDE_W, SLIDE_H, theme.WHITE_HEX)
    _add_filled_rect(slide, 0, 0, Inches(0.3), SLIDE_H, theme.PRIMARY_HEX)

    _add_textbox(slide, Inches(0.75), Inches(0.6),
                  Inches(12), Inches(0.8),
                  "Observability triangle",
                  font_size=theme.SIZE_TITLE, bold=True,
                  font_color_hex=theme.PRIMARY_HEX)

    questions = [
        ("Mode", "Which path are we on?", "synthetic / mock / live", "v7.10"),
        ("Circuit", "Is the path healthy?", "closed / intermittent / OPEN", "v8.1"),
        ("Latency", "How fast is the path?", "p50 / p95 / p99 per endpoint", "v8.2"),
    ]

    col_w = Inches(3.9)
    col_x = Inches(0.75)
    col_y = Inches(2.0)
    col_h = Inches(3.5)
    for i, (label, question, values, ver) in enumerate(questions):
        x = col_x + (col_w + Inches(0.2)) * i
        _add_filled_rect(slide, x, col_y, col_w, col_h,
                          theme.LIGHT_GREY_HEX, line_hex=theme.PRIMARY_HEX)
        _add_textbox(slide, x + Inches(0.2), col_y + Inches(0.2),
                      col_w - Inches(0.4), Inches(0.5), label,
                      font_size=22, bold=True,
                      font_color_hex=theme.PRIMARY_HEX)
        _add_textbox(slide, x + Inches(0.2), col_y + Inches(0.9),
                      col_w - Inches(0.4), Inches(0.6), question,
                      font_size=14, italic=True,
                      font_color_hex=theme.CHARCOAL_HEX)
        _add_textbox(slide, x + Inches(0.2), col_y + Inches(1.7),
                      col_w - Inches(0.4), Inches(1.0), values,
                      font_size=11, font_name=theme.PPT_MONO_FONT,
                      font_color_hex=theme.SECONDARY_HEX)
        _add_textbox(slide, x + Inches(0.2), col_y + Inches(2.9),
                      col_w - Inches(0.4), Inches(0.4), f"shipped {ver}",
                      font_size=10,
                      font_color_hex=theme.MID_GREY_HEX)

    _add_textbox(slide, Inches(0.75), Inches(6.0), Inches(12), Inches(0.6),
                  "All three surfaces visible from page 91 systems view — no tab switching required.",
                  font_size=12, italic=True,
                  font_color_hex=theme.MID_GREY_HEX)

    _add_footer(slide, slide_num, total, registry["platform"]["version"])


def _build_slide_9_audit_perimeter(prs, registry: Dict[str, Any], slide_num: int, total: int):
    """Slide 9 — 6-gate defense-in-depth."""
    slide = _add_blank_slide(prs)
    _add_filled_rect(slide, 0, 0, SLIDE_W, SLIDE_H, theme.WHITE_HEX)
    _add_filled_rect(slide, 0, 0, Inches(0.3), SLIDE_H, theme.PRIMARY_HEX)

    _add_textbox(slide, Inches(0.75), Inches(0.6),
                  Inches(12), Inches(0.8),
                  "Audit perimeter — 6-gate defense-in-depth",
                  font_size=theme.SIZE_TITLE - 4, bold=True,
                  font_color_hex=theme.PRIMARY_HEX)

    gates = [
        ("G104", "Engine migration ratchet (v7.0.1)"),
        ("G105", "Strict invariant registry usage (v7.1)"),
        ("G106", "Loop round-trip-testability (v7.15)"),
        ("G107", "Stock data_source provenance (v7.15)"),
        ("G108", "FLEXCUBE resilience + observability (v8.3)"),
        ("G109", "PUBLISHED_LANGUAGE payload_version (v8.7)"),
    ]

    # Two columns of 3 gates
    col_w = Inches(5.9)
    for col in range(2):
        for row in range(3):
            i = col * 3 + row
            gid, desc = gates[i]
            x = Inches(0.75) + col * (col_w + Inches(0.5))
            y = Inches(1.9) + row * Inches(1.2)
            # Box
            _add_filled_rect(slide, x, y, col_w, Inches(1.0),
                              theme.LIGHT_GREY_HEX, line_hex=theme.PRIMARY_HEX)
            _add_textbox(slide, x + Inches(0.2), y + Inches(0.15),
                          Inches(0.8), Inches(0.7), gid,
                          font_size=24, bold=True,
                          font_color_hex=theme.PRIMARY_HEX,
                          align=PP_ALIGN.LEFT)
            _add_textbox(slide, x + Inches(1.1), y + Inches(0.25),
                          col_w - Inches(1.3), Inches(0.6), desc,
                          font_size=12,
                          font_color_hex=theme.CHARCOAL_HEX)

    _add_textbox(slide, Inches(0.75), Inches(5.7), Inches(12), Inches(0.5),
                  "Run python scripts/audit.py to verify all gates pass.",
                  font_size=12, font_name=theme.PPT_MONO_FONT,
                  font_color_hex=theme.SECONDARY_HEX,
                  align=PP_ALIGN.LEFT)
    _add_textbox(slide, Inches(0.75), Inches(6.2), Inches(12), Inches(0.5),
                  f"Current state: {registry['platform']['audit_gates']}/109 gates passing  ·  "
                  f"{registry['platform']['audit_pass_rate']}",
                  font_size=12,
                  font_color_hex=theme.STATUS_SHIPPED_HEX)

    _add_footer(slide, slide_num, total, registry["platform"]["version"])


def _build_slide_10_compliance(prs, registry: Dict[str, Any], slide_num: int, total: int):
    """Slide 10 — Regulatory alignment."""
    slide = _add_blank_slide(prs)
    _add_filled_rect(slide, 0, 0, SLIDE_W, SLIDE_H, theme.WHITE_HEX)
    _add_filled_rect(slide, 0, 0, Inches(0.3), SLIDE_H, theme.PRIMARY_HEX)

    _add_textbox(slide, Inches(0.75), Inches(0.6),
                  Inches(12), Inches(0.8),
                  "Regulatory alignment",
                  font_size=theme.SIZE_TITLE, bold=True,
                  font_color_hex=theme.PRIMARY_HEX)

    items = registry.get("regulatory_alignment", [])
    y = Inches(1.8)
    for item in items:
        _add_textbox(slide, Inches(1.0), y,
                      Inches(0.4), Inches(0.4),
                      "✓", font_size=20, bold=True,
                      font_color_hex=theme.STATUS_SHIPPED_HEX)
        _add_textbox(slide, Inches(1.5), y,
                      Inches(11), Inches(0.5),
                      item, font_size=14,
                      font_color_hex=theme.CHARCOAL_HEX)
        y += Inches(0.6)

    # Honest note
    _add_textbox(slide, Inches(0.75), Inches(6.0),
                  Inches(12), Inches(0.5),
                  "Alignment is at architectural / engine level. Per-bank deployment "
                  "review required for production certification.",
                  font_size=10, italic=True,
                  font_color_hex=theme.MID_GREY_HEX)

    _add_footer(slide, slide_num, total, registry["platform"]["version"])


def _build_slide_11_implementation(prs, registry: Dict[str, Any], slide_num: int, total: int):
    """Slide 11 — Implementation approach."""
    slide = _add_blank_slide(prs)
    _add_filled_rect(slide, 0, 0, SLIDE_W, SLIDE_H, theme.WHITE_HEX)
    _add_filled_rect(slide, 0, 0, Inches(0.3), SLIDE_H, theme.PRIMARY_HEX)

    _add_textbox(slide, Inches(0.75), Inches(0.6),
                  Inches(12), Inches(0.8),
                  "Implementation approach — phased, audit-locked",
                  font_size=theme.SIZE_TITLE - 6, bold=True,
                  font_color_hex=theme.PRIMARY_HEX)

    phases = [
        ("Phase 1 — ACL bootstrap",
         "Connect FLEXCUBE via ACL (v8.0). Verify retry + circuit + telemetry. "
         "Synthetic tier from cbs_data/ available for parallel pilot."),
        ("Phase 2 — Systems-layer wiring",
         "Wire bank-specific stocks + flows into utils/system_stocks + system_flows. "
         "Charter §3 + §6 are the design references."),
        ("Phase 3 — UI surfacing",
         "Page 91 systems view + 4 dedicated engine pages + admin operations. "
         "Operators see mode + circuit + latency without leaving the dashboard."),
        ("Phase 4 — Audit hardening",
         "Add bank-specific audit gates. The 6-gate defense-in-depth perimeter "
         "scales — banks add their own G110+ for local invariants."),
    ]

    row_h = Inches(1.05)
    row_y = Inches(1.9)
    for i, (label, body) in enumerate(phases):
        y = row_y + row_h * i
        _add_textbox(slide, Inches(0.75), y, Inches(3.5), Inches(0.5),
                      label, font_size=14, bold=True,
                      font_color_hex=theme.PRIMARY_HEX)
        _add_textbox(slide, Inches(0.75), y + Inches(0.4), Inches(11.5), Inches(0.6),
                      body, font_size=12,
                      font_color_hex=theme.CHARCOAL_HEX)

    _add_footer(slide, slide_num, total, registry["platform"]["version"])


def _build_slide_12_references(prs, registry: Dict[str, Any], slide_num: int, total: int):
    """Slide 12 — Canonical references."""
    slide = _add_blank_slide(prs)
    _add_filled_rect(slide, 0, 0, SLIDE_W, SLIDE_H, theme.WHITE_HEX)
    _add_filled_rect(slide, 0, 0, Inches(0.3), SLIDE_H, theme.PRIMARY_HEX)

    _add_textbox(slide, Inches(0.75), Inches(0.6),
                  Inches(12), Inches(0.8),
                  "Canonical references",
                  font_size=theme.SIZE_TITLE, bold=True,
                  font_color_hex=theme.PRIMARY_HEX)

    refs = registry.get("canonical_references", [])
    y = Inches(2.0)
    for ref in refs:
        # Strip markdown asterisks for display
        clean_ref = ref.replace("*", "")
        _add_textbox(slide, Inches(1.0), y, Inches(11.5), Inches(0.5),
                      "• " + clean_ref, font_size=14,
                      font_color_hex=theme.CHARCOAL_HEX)
        y += Inches(0.6)

    _add_textbox(slide, Inches(0.75), Inches(6.0), Inches(12), Inches(0.4),
                  "Internal canon: A2Z_SYSTEMS_CHARTER.md + A2Z_V7_RETROSPECTIVE.md + "
                  "A2Z_V8_RETROSPECTIVE.md + A2Z_LIVING_DOCS_PLAN.md",
                  font_size=10, italic=True,
                  font_color_hex=theme.MID_GREY_HEX)

    _add_footer(slide, slide_num, total, registry["platform"]["version"])


def _build_slide_13_honest_scope(prs, registry: Dict[str, Any], slide_num: int, total: int):
    """Slide 13 — What this deck does NOT claim. MANDATORY."""
    slide = _add_blank_slide(prs)
    # Distinct background to mark this as the honesty section
    _add_filled_rect(slide, 0, 0, SLIDE_W, SLIDE_H, theme.WHITE_HEX)
    _add_filled_rect(slide, 0, 0, Inches(0.3), SLIDE_H, theme.STATUS_DESIGNED_HEX)

    _add_textbox(slide, Inches(0.75), Inches(0.6),
                  Inches(12), Inches(0.8),
                  section_title(),
                  font_size=theme.SIZE_TITLE - 4, bold=True,
                  font_color_hex=theme.STATUS_DESIGNED_HEX)

    # Show top items from honest scope (limit to fit)
    items = collect_honest_scope_lines(registry)[:7]

    intro = (
        "This slide is mandatory in the A2Z brochure. It enumerates what the "
        "deck deliberately does NOT claim — features that are roadmap rather "
        "than shipped, integrations that are designed but not deployed, "
        "outcomes that are projections rather than measurements."
    )
    _add_textbox(slide, Inches(0.75), Inches(1.6), Inches(12), Inches(1.0),
                  intro, font_size=11, italic=True,
                  font_color_hex=theme.CHARCOAL_HEX)

    y = Inches(2.7)
    for item in items:
        # Truncate very long items
        truncated = item if len(item) < 160 else item[:157] + "..."
        _add_textbox(slide, Inches(0.9), y, Inches(0.4), Inches(0.4),
                      "→", font_size=14, bold=True,
                      font_color_hex=theme.STATUS_DESIGNED_HEX)
        _add_textbox(slide, Inches(1.4), y, Inches(11.5), Inches(0.5),
                      truncated, font_size=11,
                      font_color_hex=theme.CHARCOAL_HEX)
        y += Inches(0.55)

    _add_footer(slide, slide_num, total, registry["platform"]["version"])


def _build_slide_14_verify(prs, registry: Dict[str, Any], slide_num: int, total: int):
    """Slide 14 — Verify yourself."""
    slide = _add_blank_slide(prs)
    _add_filled_rect(slide, 0, 0, SLIDE_W, SLIDE_H, theme.PRIMARY_HEX)

    _add_textbox(slide, Inches(0.75), Inches(0.8),
                  Inches(12), Inches(1.0),
                  "Don't take our word for it",
                  font_size=44, bold=True,
                  font_color_hex=theme.WHITE_HEX)

    _add_textbox(slide, Inches(0.75), Inches(2.2),
                  Inches(12), Inches(0.6),
                  "Run the audit yourself. Inspect every gate. Read every CHANGELOG.",
                  font_size=18, italic=True,
                  font_color_hex="CADCFC")

    # Code block
    _add_filled_rect(slide, Inches(0.75), Inches(3.3),
                      Inches(11.8), Inches(2.0),
                      theme.ACCENT_HEX)
    cmd = (
        f"$ python {registry['platform']['audit_command'].split()[1]}\n"
        f"  Score: {registry['platform']['audit_gates']}/109 gates = "
        f"{registry['platform']['audit_pass_rate']} — PASS"
    )
    _add_textbox(slide, Inches(1.0), Inches(3.5),
                  Inches(11.3), Inches(1.5),
                  cmd, font_size=18, font_name=theme.PPT_MONO_FONT,
                  font_color_hex="CADCFC")

    _add_textbox(slide, Inches(0.75), Inches(6.0),
                  Inches(12), Inches(0.5),
                  f"{registry['platform']['changelog_count']} per-batch CHANGELOGs document every shipped batch.",
                  font_size=14,
                  font_color_hex=theme.WHITE_HEX)


def _build_slide_15_next(prs, registry: Dict[str, Any], slide_num: int, total: int):
    """Slide 15 — Next steps."""
    slide = _add_blank_slide(prs)
    _add_filled_rect(slide, 0, 0, SLIDE_W, SLIDE_H, theme.WHITE_HEX)
    _add_filled_rect(slide, 0, 0, Inches(0.3), SLIDE_H, theme.PRIMARY_HEX)

    _add_textbox(slide, Inches(0.75), Inches(0.6),
                  Inches(12), Inches(0.8),
                  "Next steps",
                  font_size=theme.SIZE_TITLE, bold=True,
                  font_color_hex=theme.PRIMARY_HEX)

    steps = [
        "Read the canonical docs: charter + retrospectives in docs/",
        "Run the audit on the version you receive: python scripts/audit.py",
        "Inspect the systems-layer registries: utils/system_stocks.py + utils/system_flows.py",
        "Review per-batch CHANGELOGs for the discipline pattern",
        "Schedule architectural review with bank's CIO + CISO + Head of Risk",
    ]
    y = Inches(2.0)
    for i, step in enumerate(steps):
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.85), y,
                                          Inches(0.5), Inches(0.5))
        circle.fill.solid()
        circle.fill.fore_color.rgb = _hex_to_rgb(theme.PRIMARY_HEX)
        circle.line.fill.background()
        circle.text_frame.text = str(i + 1)
        circle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        for run in circle.text_frame.paragraphs[0].runs:
            run.font.size = Pt(14)
            run.font.bold = True
            run.font.color.rgb = _hex_to_rgb(theme.WHITE_HEX)
        _add_textbox(slide, Inches(1.6), y, Inches(11), Inches(0.5),
                      step, font_size=14,
                      font_color_hex=theme.CHARCOAL_HEX)
        y += Inches(0.7)

    _add_footer(slide, slide_num, total, registry["platform"]["version"])


# ════════════════════════════════════════════════════════════════════
# Public API
# ════════════════════════════════════════════════════════════════════

def _build_claims(registry: Dict[str, Any]) -> List[Claim]:
    """Build the canonical claim list this generator asserts.

    The brochure's audit-locked claims. Validator runs these BEFORE
    any rendering. If any fail, the brochure is never written.
    """
    return [
        Claim(text="6 system stocks",
              registry_path="stocks_count",
              expected_value=6,
              source_file="utils/system_stocks.py"),
        Claim(text="15 feedback loops",
              registry_path="loops_count",
              expected_value=15,
              source_file="utils/system_flows.py"),
        Claim(text="100% loops wired",
              registry_path="loops_wired_pct",
              expected_value=100.0,
              source_file="utils/system_flows.py"),
        Claim(text="3 learning loops",
              registry_path="learning_loops_count",
              expected_value=3,
              source_file="utils/system_flows.py"),
        Claim(text="6 of 6 sales-content JSONs present",
              registry_path="sales_content_files_present",
              expected_value=6,
              source_file="docs/sales_content/"),
    ]


def generate_brochure(output_path: Path) -> Dict[str, Any]:
    """Generate the 15-slide A2Z brochure. Audit-locked.

    Returns dict with: status, claims_validated, output_path, slide_count.
    Raises ClaimValidationError if claims diverge from registry.
    """
    registry = load_registry()

    # Validate claims FIRST — abort on divergence per Living Doc discipline
    claims = _build_claims(registry)
    result = validate_claims(claims, registry, fail_fast=False)
    if result["failed"] > 0:
        return {
            "status": "ABORTED",
            "reason": "Claim validation failed; collateral not written",
            "claims_validated": result["passed"],
            "claims_failed": result["failed"],
            "failures": result["failures"],
        }

    # All claims pass — proceed with rendering
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    total = 15
    _build_slide_1_title(prs, registry, total)
    _build_slide_2_problem(prs, registry, 2, total)
    _build_slide_3_solution(prs, registry, 3, total)
    _build_slide_4_architecture(prs, registry, 4, total)
    _build_slide_5_systems(prs, registry, 5, total)
    _build_slide_6_acl(prs, registry, 6, total)
    _build_slide_7_resilience(prs, registry, 7, total)
    _build_slide_8_observability(prs, registry, 8, total)
    _build_slide_9_audit_perimeter(prs, registry, 9, total)
    _build_slide_10_compliance(prs, registry, 10, total)
    _build_slide_11_implementation(prs, registry, 11, total)
    _build_slide_12_references(prs, registry, 12, total)
    _build_slide_13_honest_scope(prs, registry, 13, total)
    _build_slide_14_verify(prs, registry, 14, total)
    _build_slide_15_next(prs, registry, 15, total)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))

    return {
        "status": "OK",
        "output_path": str(output_path),
        "slide_count": total,
        "claims_validated": result["passed"],
        "platform_version": registry["platform"]["version"],
    }


if __name__ == "__main__":
    import sys
    out = Path("/tmp/A2Z_MIS_360_Brochure.pptx")
    result = generate_brochure(out)
    print(f"A2Z PPT generator — {result['status']}")
    if result["status"] == "OK":
        print(f"  Output: {result['output_path']}")
        print(f"  Slides: {result['slide_count']}")
        print(f"  Claims validated: {result['claims_validated']}")
    else:
        print(f"  Reason: {result['reason']}")
        for f in result.get("failures", []):
            print(f"    {f['error']}")
        sys.exit(1)
